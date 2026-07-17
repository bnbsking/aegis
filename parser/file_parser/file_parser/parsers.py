import base64
import json
import os
import subprocess
from tempfile import TemporaryDirectory
from typing import Dict, List, Literal

import pandas as pd
import pymupdf
import extract_msg
from pdf2image import convert_from_path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from file_parser.ocr_client import request_ocr


class BaseParser:
    def __init__(self, input_path: str, output_dir: str | None = None, **kwargs):
        self.input_path = input_path
        self.output_dir = output_dir
        self.kwargs = kwargs

    def write(self, out: str | Dict | List) -> None:
        if self.output_dir:
            os.makedirs(self.output_dir, exist_ok=True)
            output_path_prefix = os.path.join(
                self.output_dir,
                f"{os.path.splitext(os.path.basename(self.input_path))[0]}"
            )
            if isinstance(out, str):
                with open(f"{output_path_prefix}.txt", "w", encoding="utf-8") as f:
                    f.write(out)
            elif isinstance(out, dict) or isinstance(out, list):
                with open(f"{output_path_prefix}.json", "w", encoding="utf-8") as f:
                    json.dump(out, f, ensure_ascii=False, indent=4)
            else:
                raise ValueError(f"Unsupported output type: {type(out)}")

    def run(self) -> str | Dict | List:
        raise NotImplementedError("Subclasses must implement the run method.")


class PDFParser(BaseParser):
    def __init__(
            self,
            input_path: str,
            output_dir: str | None = None,
            ocr_mode: Literal['text', 'hybrid', 'ocr'] = 'ocr',
            pdf_url: str = "http://172.24.37.129:8002/run_pdf_list",
            img_url: str = "http://172.24.37.129:8002/run_img_list",
            to_str: bool = False
        ):
        super().__init__(input_path, output_dir)
        self.ocr_mode = ocr_mode
        self.pdf_url = pdf_url
        self.img_url = img_url
        self.to_str = to_str

    @staticmethod
    def get_need_ocr_pages(text_list: List[str], min_chars: int = 50, chars_ratio_threshold: float = 0.3) -> List[int]:
        ocr_page_list = []
        for i, text in enumerate(text_list):
            if len(text.strip()) <= min_chars:
                ocr_page_list.append(i)
            elif sum(c.isalpha() for c in text) < len(text) * chars_ratio_threshold:
                ocr_page_list.append(i)
        return ocr_page_list

    @staticmethod
    def get_imgs_from_pdf(pdf_path: str, pages: List[int]) -> List[bytes]:
        images = convert_from_path(pdf_path, dpi=300)  # List[PIL]
        img_list = []
        with TemporaryDirectory() as tmp_dir:
            for page_idx in pages:
                img_path = f"{tmp_dir}/page_{page_idx}.jpg"
                images[page_idx].save(img_path, "JPEG")
                with open(img_path, "rb") as f:
                    img_list.append(f.read())
        return img_list

    @staticmethod
    def hybrid_merge(text_list: List[str], ocr_result: List[str]) -> List[str]:
        """
        Here expect ocr_result is List[str], where each element is the text of each page
        If ocr_server has different output, just override this function in subclass
        """
        merge_list = []
        for i in range(max(len(text_list), len(ocr_result))):
            merge_list.append(
                (text_list[i] if i < len(text_list) else "")\
                + (ocr_result[i] if i < len(ocr_result) else "")
            )
        return merge_list
    
    @staticmethod
    def get_text_list(pdf_path: str) -> List[str]:
        with pymupdf.open(pdf_path) as doc:
            text_list = [page.get_text("text") for page in doc]
        return text_list
        
    def run(self) -> str | Dict | List:
        if self.ocr_mode == 'text':
            text_list = self.get_text_list(self.input_path)
            out = "\n".join(text_list) if self.to_str else text_list
        
        elif self.ocr_mode == 'hybrid':
            text_list = self.get_text_list(self.input_path)
            need_ocr_pages = self.get_need_ocr_pages(text_list)
            if need_ocr_pages:
                img_list = self.get_imgs_from_pdf(self.input_path, need_ocr_pages)
                ocr_result = request_ocr(img_list, self.img_url)
                out: List[str] = self.hybrid_merge(text_list, ocr_result)
            else:
                out = text_list
            out = "\n".join(out) if self.to_str else out
            
        elif self.ocr_mode == 'ocr':
            assert self.pdf_url, "pdf_url must be provided for ocr mode"
            out = request_ocr([open(self.input_path, "rb").read()], self.pdf_url)
        
        else:
            raise ValueError(f"Unsupported ocr_mode: {self.ocr_mode}")
        
        self.write(out)
        return out


class WordParser(BaseParser):
    """
    This class change word to pdf and apply PDFParser
    """
    def __init__(
            self,
            input_path: str,
            output_dir: str | None = None,
            ocr_mode: Literal['text', 'hybrid', 'ocr'] = 'hybrid',
            pdf_url: str = "http://172.24.37.129:8002/run_pdf_list",
            img_url: str = "http://172.24.37.129:8002/run_img_list",
            to_str: bool = False
        ):
        super().__init__(input_path, output_dir)
        self.ocr_mode = ocr_mode
        self.pdf_url = pdf_url
        self.img_url = img_url
        self.to_str = to_str

    @staticmethod
    def word2pdf_with_path(input_path: str, tmp_dir: str) -> str:
        cmd = [
            "libreoffice",
            "--headless",
            "--convert-to", "pdf",
            str(input_path),
            "--outdir", str(tmp_dir)
        ]
        subprocess.run(cmd, check=True)
        basename, _ = os.path.splitext(os.path.basename(input_path))
        output_path = os.path.join(tmp_dir,  f"{basename}.pdf")
        return output_path

    def run(self) -> str | Dict | List:
        with TemporaryDirectory() as tmp_dir:
            pdf_path = self.word2pdf_with_path(self.input_path, tmp_dir)
            pdf_parser = PDFParser(
                pdf_path,
                self.output_dir,
                self.ocr_mode,
                self.pdf_url,
                self.img_url,
                self.to_str
            )
            out = pdf_parser.run()
        return out


class ExcelParser(BaseParser):
    def __init__(
            self,
            input_path: str,
            output_dir: str | None = None,
            to_str: bool = False
        ):
        super().__init__(input_path, output_dir)
        self.to_str = to_str

    def run(self) -> str | Dict | List:
        if self.input_path.lower().endswith("csv"):
            df = pd.read_csv(self.input_path, encoding="utf-8")
            if self.to_str:
                out = df.to_string(index=False)
            else:
                out = df.to_dict(orient="records")
        else:
            text_list = []
            xls = pd.ExcelFile(self.input_path)
            for sheet_name in xls.sheet_names:
                df = pd.read_excel(self.input_path, sheet_name=sheet_name)
                if self.to_str:
                    text_list.append(df.to_string(index=False))
                else:
                    text_list.append(df.to_dict(orient="records"))
            out = "\n".join(text_list) if self.to_str else text_list
        self.write(out)
        return out


class DummyParser(BaseParser):
    def run(self) -> str:
        with open(self.input_path, "r", encoding="utf-8") as f:
            out = f.read()
        self.write(out)
        return out


class ImageParser(BaseParser):
    def __init__(
            self,
            input_path: str,
            output_dir: str | None = None,
            img_url: str = "http://172.24.37.129:8002/run_img_list"
        ):
        super().__init__(input_path, output_dir)
        self.img_url = img_url

    def run(self) -> str:
        out = request_ocr([open(self.input_path, "rb").read()], self.img_url)
        self.write(out)
        return out


class MsgParser(BaseParser):
    """Parse an Outlook .msg file into subject/sender/body/images."""

    def run(self) -> Dict:
        with extract_msg.openMsg(self.input_path) as msg:
            images = []
            for att in (msg.attachments or []):
                mimetype = getattr(att, "mimetype", "") or ""
                if not mimetype.startswith("image/"):
                    continue
                raw = getattr(att, "data", None)
                if not raw:
                    continue
                filename = getattr(att, "longFilename", "") or getattr(att, "shortFilename", "") or "image"
                images.append({
                    "filename": filename,
                    "mimetype": mimetype,
                    "data": base64.b64encode(raw).decode("ascii"),
                })
            out = {
                "subject": (msg.subject or "").strip(),
                "sender": (msg.sender or "").strip(),
                "to": (msg.to or "").strip(),
                "cc": (msg.cc or "").strip(),
                "body": (msg.body or "").strip(),
                "images": images,
            }

        self.write(out)
        return out


class PPTParser(BaseParser):
    """Parse a PowerPoint file into per-slide text, tables, and images (no OCR)."""

    @staticmethod
    def ppt2pptx_with_path(input_path: str, tmp_dir: str) -> str:
        cmd = [
            "libreoffice",
            "--headless",
            "--convert-to", "pptx",
            str(input_path),
            "--outdir", str(tmp_dir)
        ]
        subprocess.run(cmd, check=True)
        basename, _ = os.path.splitext(os.path.basename(input_path))
        output_path = os.path.join(tmp_dir, f"{basename}.pptx")
        return output_path

    @staticmethod
    def flatten_shapes(shapes):
        for shape in shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                yield from PPTParser.flatten_shapes(shape.shapes)
            else:
                yield shape

    def run(self) -> List[Dict]:
        with TemporaryDirectory() as tmp_dir:
            input_path = self.input_path
            if input_path.lower().endswith(".ppt"):
                input_path = self.ppt2pptx_with_path(input_path, tmp_dir)

            prs = Presentation(input_path)
            slides = []
            for slide in prs.slides:
                texts = []
                tables = []
                images = []
                shapes = self.flatten_shapes(slide.shapes)
                sorted_shapes = sorted(shapes, key=lambda s: (s.top or 0, s.left or 0))
                for shape in sorted_shapes:
                    if shape.has_text_frame and shape.text_frame.text.strip():
                        texts.append(shape.text_frame.text)
                    if shape.has_table:
                        tables.append([[cell.text for cell in row.cells] for row in shape.table.rows])
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        image = shape.image
                        images.append({
                            "filename": f"image_{len(images)}.{image.ext}",
                            "mimetype": f"image/{image.ext}",
                            "data": base64.b64encode(image.blob).decode("ascii"),
                        })
                slides.append({
                    "text": "\n".join(texts),
                    "tables": tables,
                    "images": images,
                })

        self.write(slides)
        return slides
